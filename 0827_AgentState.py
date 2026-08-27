# pip install PyPDF2    
# pip install python-docx
# pip install python-pptx
# pip install torchvision

import streamlit as st
import base64
import uuid
import pandas as pd
import pypdf
from docx import Document
from pptx import Presentation
from typing import Annotated, TypedDict, Union
from langchain_ollama import ChatOllama
from langchain_core.messages import HumanMessage, AIMessage, BaseMessage
from langgraph.graph import StateGraph, END
from langgraph.checkpoint.memory import MemorySaver

# --- 1. 모델 설정 ---
ROUTER_MODEL = "mistral:latest"
GENERATOR_MODEL = "gemma3:latest"

router_llm = ChatOllama(model=ROUTER_MODEL, temperature=0)
generator_llm = ChatOllama(model=GENERATOR_MODEL, temperature=0.7)

# --- 2. LangGraph 상태 정의 ---
class AgentState(TypedDict):
    messages: Annotated[list[BaseMessage], "대화 내역"]
    category: str
    image_data: Union[str, None]
    doc_text: Union[str, None]

# --- 3. 에이전트 노드 함수 ---
def router_node(state: AgentState):
    last_message = state['messages'][-1].content
    prompt = f"분류 전문가로서 다음 요청을 [광고, 유튜브, SNS] 중 하나로 분류하세요. 단어만 출력: {last_message}"
    category = router_llm.invoke(prompt).content.strip()
    return {"category": category}

def generator_node(state: AgentState):
    category = state.get("category", "일반")
    messages = state['messages']
    image_data = state.get("image_data")
    doc_text = state.get("doc_text")
    
    role_map = {
        "광고": "카피라이터. 제품 강점 부각 및 구매 유도 전문.",
        "유튜브": "영상 작가. [섹션명], 🎬화면, 🎙️내레이션 구분의 상세 대본 작성 전문.",
        "SNS": "인플루언서. 트렌디한 문체, 풍부한 이모지, 해시태그 활용 전문."
    }
    persona = role_map.get(category, "전문가")
    system_msg = f"당신은 전문 {persona}입니다. 모든 답변은 한국어로 작성하세요."
    
    if doc_text:
        system_msg += f"\n\n[참고 문서 내용]\n{doc_text}"
    
    if image_data:
        last_text = messages[-1].content
        content = [
            {"type": "text", "text": f"{system_msg}\n\n사용자 요청: {last_text}"},
            {"type": "image_url", "image_url": f"data:image/jpeg;base64,{image_data}"}
        ]
        final_input = [HumanMessage(content=content)]
    else:
        final_input = [AIMessage(content=system_msg)] + messages

    response = generator_llm.invoke(final_input)
    return {"messages": [response]}

# --- 4. 워크플로우 빌드 ---
workflow = StateGraph(AgentState)
workflow.add_node("router", router_node)
workflow.add_node("generator", generator_node)
workflow.set_entry_point("router")
workflow.add_edge("router", "generator")
workflow.add_edge("generator", END)

memory = MemorySaver()
app = workflow.compile(checkpointer=memory)

# --- 그래프 시각화 이미지 생성 및 저장 ---
try:
    png_data = app.get_graph().draw_mermaid_png()
    with open("./images2/langGraph04_streamlit.png", "wb") as f:
        f.write(png_data)
except Exception as e:
    import sys
    print(f"그래프 시각화 이미지 저장 중 오류 발생: {e}", file=sys.stderr)

# --- 5. 유틸리티 함수 ---
def get_document_text(uploaded_file):
    name = uploaded_file.name
    ext = name.split('.')[-1].lower()
    text = ""
    try:
        if ext == 'pdf':
            reader = pypdf.PdfReader(uploaded_file)
            for page in reader.pages: text += page.extract_text()
        elif ext == 'docx':
            doc = Document(uploaded_file)
            for para in doc.paragraphs: text += para.text + "\n"
        elif ext == 'pptx':
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
        elif ext in ['csv', 'xlsx']:
            df = pd.read_csv(uploaded_file) if ext == 'csv' else pd.read_excel(uploaded_file)
            text = df.to_string()
        elif ext == 'txt':
            text = uploaded_file.read().decode("utf-8")
    except Exception as e:
        text = f"파일 읽기 오류: {str(e)}"
    return text

def summarize_title(text):
    summary_prompt = f"다음 내용을 10자 이내 한국어 제목으로 요약해줘: {text}"
    try: return router_llm.invoke(summary_prompt).content.strip().replace('"', '')
    except: return "새 대화"

# --- 6. UI 레이아웃 ---
st.set_page_config(page_title="AI 통합 콘텐츠 제작소", layout="wide")

# 세션 상태 초기화
if "chat_sessions" not in st.session_state: st.session_state.chat_sessions = {}
if "show_all_history" not in st.session_state: st.session_state.show_all_history = False
if "current_thread_id" not in st.session_state:
    new_id = str(uuid.uuid4())
    st.session_state.current_thread_id = new_id
    st.session_state.chat_sessions[new_id] = {"title": "새 대화", "messages": []}

with st.sidebar:
    st.title("🛠️ 제어 센터")
    if st.button("➕ 새 채팅 시작", use_container_width=True):
        new_id = str(uuid.uuid4())
        st.session_state.current_thread_id = new_id
        st.session_state.chat_sessions[new_id] = {"title": "새 대화", "messages": []}
        st.session_state.temp_hint = None
        st.rerun()

    st.divider()
    st.subheader("🎯 컨셉 선택")
    c1, c2, c3 = st.columns(3)
    if c1.button("📺", help="유튜브 스크립트 모드"):
        st.session_state.temp_hint = "유튜브 스크립트 작성해줘: "
        st.toast("📺 유튜브 모드 활성화", icon="✅")
    if c2.button("📢", help="광고 카피 모드"):
        st.session_state.temp_hint = "광고 카피 작성해줘: "
        st.toast("📢 광고 모드 활성화", icon="✅")
    if c3.button("📱", help="SNS 포스팅 모드"):
        st.session_state.temp_hint = "SNS 포스팅 작성해줘: "
        st.toast("📱 SNS 모드 활성화", icon="✅")
    
    if st.session_state.get("temp_hint"):
        mode_name = st.session_state.temp_hint.split(" ")[0]
        st.success(f"현재 활성: **{mode_name}**")
    
    st.divider()
    st.subheader("📜 대화 히스토리")
    
    # [수정된 히스토리 표시 로직]
    session_keys = list(st.session_state.chat_sessions.keys())[::-1] # 최신순
    display_limit = 5
    
    # 보여줄 리스트 결정
    if st.session_state.show_all_history or len(session_keys) <= display_limit:
        visible_sessions = session_keys
    else:
        visible_sessions = session_keys[:display_limit]

    for tid in visible_sessions:
        col_list, col_del = st.columns([0.8, 0.2])
        with col_list:
            label = st.session_state.chat_sessions[tid]["title"]
            is_cur = (tid == st.session_state.current_thread_id)
            if st.button(f"{'📍 ' if is_cur else ''}{label[:10]}", key=f"btn_{tid}", use_container_width=True):
                st.session_state.current_thread_id = tid
                st.rerun()
        with col_del:
            if st.button("🗑️", key=f"del_{tid}"):
                del st.session_state.chat_sessions[tid]
                if tid == st.session_state.current_thread_id:
                    st.session_state.current_thread_id = None
                st.rerun()

    # 펼치기/접기 버튼 제어
    if len(session_keys) > display_limit:
        if st.session_state.show_all_history:
            if st.button("🔼 접기", use_container_width=True):
                st.session_state.show_all_history = False
                st.rerun()
        else:
            if st.button(f"🔽 더 보기 ({len(session_keys)-display_limit}개 더 있음)", use_container_width=True):
                st.session_state.show_all_history = True
                st.rerun()

    st.divider()
    uploaded_file = st.file_uploader("📎 파일 분석 (이미지/영상/문서)", type=['jpg', 'png', 'mp4', 'pdf', 'docx', 'pptx', 'csv', 'xlsx', 'txt'])


# 메인 화면
st.title("🚀 Multi-Agent 콘텐츠 생성기")

if st.session_state.current_thread_id is None or st.session_state.current_thread_id not in st.session_state.chat_sessions:
    new_id = str(uuid.uuid4())
    st.session_state.current_thread_id = new_id
    st.session_state.chat_sessions[new_id] = {"title": "새 대화", "messages": []}

current_session = st.session_state.chat_sessions[st.session_state.current_thread_id]

# 메시지 출력 (질문 먼저 보이게 구성)
for msg in current_session["messages"]:
    with st.chat_message(msg["role"]):
        if msg["role"] == "user" and "file_info" in msg:
            f = msg["file_info"]
            if f["type"] == "image": st.image(f["data"], width=400)
            elif f["type"] == "video": st.video(f["data"])
            else: st.info(f"📄 파일명: {f['name']} ({f['ext'].upper()})")
        st.markdown(msg["content"])

if st.session_state.get("temp_hint"):
    st.warning(f"💡 현재 **{st.session_state.temp_hint.split(' ')[0]}** 모드로 답변이 생성됩니다.")

# 입력 처리
if prompt := st.chat_input("메시지를 입력하세요..."):
    prefix = st.session_state.get("temp_hint") or ""
    full_prompt = prefix + prompt
    st.session_state.temp_hint = None
    
    user_msg_obj = {"role": "user", "content": full_prompt}
    
    if uploaded_file:
        ext = uploaded_file.name.split('.')[-1].lower()
        if ext in ['jpg', 'jpeg', 'png']:
            data = uploaded_file.read()
            user_msg_obj["file_info"] = {"type": "image", "data": data}
            user_msg_obj["b64_data"] = base64.b64encode(data).decode("utf-8")
        elif ext == 'mp4':
            user_msg_obj["file_info"] = {"type": "video", "data": uploaded_file.read()}
        else:
            text = get_document_text(uploaded_file)
            user_msg_obj["file_info"] = {"type": "doc", "name": uploaded_file.name, "ext": ext}
            user_msg_obj["doc_text"] = text

    if not current_session["messages"]:
        current_session["title"] = summarize_title(full_prompt)
    
    current_session["messages"].append(user_msg_obj)
    st.rerun()

# 답변 생성
if current_session["messages"] and current_session["messages"][-1]["role"] == "user":
    last_user_msg = current_session["messages"][-1]
    with st.chat_message("assistant"):
        with st.status("🛠️ 에이전트 협업 분석 중...", expanded=False) as status:
            config = {"configurable": {"thread_id": st.session_state.current_thread_id}}
            inputs = {
                "messages": [HumanMessage(content=last_user_msg["content"])], 
                "image_data": last_user_msg.get("b64_data"), 
                "doc_text": last_user_msg.get("doc_text")
            }
            final_state = app.invoke(inputs, config)
            response = final_state["messages"][-1].content
            status.update(label=f"✅ {final_state['category']} 생성 완료", state="complete")
        st.markdown(response)
        current_session["messages"].append({"role": "assistant", "content": response})
        st.rerun()


# ModuleNotFoundError: No module named 'docx'
# pip install python-docx

# ModuleNotFoundError: No module named 'pptx'
# pip install python-pptx 

